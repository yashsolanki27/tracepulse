"""Generate TracePulse ARD as PDF (reportlab) and PPTX (python-pptx) in docs/."""
from pathlib import Path

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from pptx import Presentation
from pptx.util import Inches, Pt

DOCS = Path(__file__).resolve().parent.parent / "docs"
DOCS.mkdir(exist_ok=True)

TITLE = "TracePulse â€” Architecture & Requirements Document"

SECTIONS = [
    ("Section 1 â€” Overview",
     "Standalone AI incident-ticket tool: RCA (root cause analysis) + similarity search. "
     "Ticket comes in â†’ AI diagnoses it â†’ system finds similar past resolved incidents â†’ "
     "engineer reviews/fixes â†’ resolution feeds back into future similarity matches. "
     "No dependency on other projects (Pulsegrid/LogPulse)."),
    ("Section 2 â€” Stack",
     "FastAPI, PostgreSQL + pgvector, Alembic, Groq API (openai/gpt-oss-120b) for RCA, "
     "sentence-transformers (all-MiniLM-L6-v2, CPU-only) for embeddings, "
     "Docker Compose (local), static API key auth. 100% free tier."),
    ("Section 3 â€” V1 Status: SHIPPED (tagged v1.0)", None),
    ("Section 4 â€” V1 Endpoints",
     "POST /tickets, GET /tickets/{id}, GET /tickets, PATCH /tickets/{id}/resolve â€” "
     "all require X-API-Key header"),
    ("Section 5 â€” V2 Roadmap (locked scope, phased)", None),
    ("Section 6 â€” Key Architectural Decisions",
     "â€¢ Model swapped from spec'd Llama 3.1 (deprecated) to openai/gpt-oss-120b\n"
     "â€¢ RCA/triage always fail-safe: null fields on timeout/error, ticket never fails to save\n"
     "â€¢ 10s timeout, retry-once on bad JSON\n"
     "â€¢ Alembic used from Phase 2 onward for all schema changes\n"
     "â€¢ Seed data generated through real API calls so embeddings/RCA are genuine, "
     "not placeholder"),
    ("Section 7 â€” Deployment",
     "Local Docker Compose only (2 containers: api, db). No cloud deploy yet."),
]

V1_TABLE = {
    "header": ["Phase", "What It Does", "Status", "Key Issue Resolved"],
    "rows": [
        ["1. Scaffold", "Docker Compose + FastAPI skeleton", "DONE", "health route cleanup"],
        ["2. Tickets table + pgvector + Alembic", "", "DONE",
         "missing CREATE EXTENSION in migration"],
        ["3. Plain CRUD + validation", "", "DONE", "clean pass"],
        ["4. API key auth", "", "DONE",
         "weak AI-generated key replaced with real entropy (openssl)"],
        ["5. RCA wiring (Groq)", "", "DONE",
         "spec'd model deprecated, swapped to gpt-oss-120b; fixed import/dependency/DB-host crash-loop"],
        ["6. Similarity search (pgvector)", "", "DONE",
         "GPU/CUDA torch bloat fixed via correct --index-url pinning"],
        ["7. Resolve endpoint", "", "DONE",
         "NULL-embedding crash in similarity query fixed"],
        ["8. Seed script", "18 realistic tickets across 6 domains, seeded via real API calls (not fake DB inserts)", "DONE",
         "stale env var causing silent 401s fixed"],
        ["9. End-to-end ship gate", "", "DONE",
         "cold restart, all 4 endpoints verified, 29 tickets survived restart intact"],
    ],
}

V2_TABLE = {
    "header": ["Phase", "Scope", "Dependency", "Status"],
    "rows": [
        ["2a", "Incident triage (priority/severity/issue_type/team) + resolve/close state machine",
         "No blockers", "IN PROGRESS"],
        ["2b", "SLA management + monitoring/escalation", "Needs background scheduler", "Not started"],
        ["2c", "Engineer assignment + notification",
         "Needs user table + notification channel", "Not started"],
        ["2d", "Unified engineer dashboard (frontend)", "Frontend stack TBD", "Not started"],
        ["2e", "Email ingestion (IMAP/Graph)", "Isolated, no dependency", "Not started"],
        ["Blocked", "Pulsegrid webhook integration",
         "Blocked on Pulsegrid's own deploy, spec ready", "No ETA"],
    ],
}

# ---------------- PDF ----------------
def make_pdf(path: Path):
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("SectionHeader", parent=styles["Heading1"], fontSize=14,
                        spaceBefore=14, spaceAfter=6, textColor=colors.HexColor("#1a3c6e"))
    body = ParagraphStyle("Body", parent=styles["BodyText"], fontSize=10, leading=14)
    title = ParagraphStyle("DocTitle", parent=styles["Title"], fontSize=20, spaceAfter=4)

    def table(spec, widths):
        data = [spec["header"]] + spec["rows"]
        t = Table(data, colWidths=widths, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1a3c6e")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#eef2f8")]),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ]))
        return t

    story = [Paragraph(TITLE, title), Spacer(1, 8)]
    for header, text in SECTIONS:
        story.append(Paragraph(header, h1))
        if header.startswith("Section 3"):
            story.append(table(V1_TABLE, [40 * mm, 55 * mm, 18 * mm, 67 * mm]))
        elif header.startswith("Section 5"):
            story.append(table(V2_TABLE, [20 * mm, 70 * mm, 55 * mm, 35 * mm]))
        elif text:
            for line in text.split("\n"):
                story.append(Paragraph(line, body))
    SimpleDocTemplate(str(path), pagesize=A4, title=TITLE).build(story)


# ---------------- PPTX ----------------
def make_pptx(path: Path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide(title, body=None, table_spec=None):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size, p.font.bold = Pt(32), True
        if body:
            bb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.6))
            tf = bb.text_frame
            tf.word_wrap = True
            for i, line in enumerate(body.split("\n")):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = line
                para.font.size = Pt(18)
        if table_spec:
            rows = len(table_spec["rows"]) + 1
            cols = len(table_spec["header"])
            shape = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.3),
                                           Inches(12.5), Inches(0.65) * rows)
            table = shape.table
            for c, h in enumerate(table_spec["header"]):
                cell = table.cell(0, c)
                cell.text = h
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(14)
            for r, row in enumerate(table_spec["rows"], start=1):
                for c, val in enumerate(row):
                    cell = table.cell(r, c)
                    cell.text = val
                    cell.text_frame.paragraphs[0].font.size = Pt(11 if rows > 6 else 14)

    # Title slide
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = TITLE
    p.font.size, p.font.bold = Pt(40), True

    for header, text in SECTIONS:
        if header.startswith("Section 3"):
            add_slide(header, table_spec=V1_TABLE)
        elif header.startswith("Section 5"):
            add_slide(header, table_spec=V2_TABLE)
        else:
            add_slide(header, body=text)

    prs.save(str(path))


if __name__ == "__main__":
    pdf, pptx = DOCS / "tracepulse_ard.pdf", DOCS / "tracepulse_ard.pptx"
    make_pdf(pdf)
    make_pptx(pptx)
    print("PDF:", pdf, pdf.stat().st_size, "bytes")
    print("PPTX:", pptx, pptx.stat().st_size, "bytes")
