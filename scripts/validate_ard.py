"""Validate the generated ARD files."""
import PyPDF2
from pptx import Presentation
from zipfile import ZipFile

r = PyPDF2.PdfReader("docs/tracepulse_ard.pdf")
text = "\n".join(p.extract_text() for p in r.pages)
print("PDF pages:", len(r.pages))
for s in ["Overview", "Stack", "SHIPPED", "Endpoints", "Roadmap", "Architectural", "Deployment", "gpt-oss-120b", "pgvector"]:
    assert s in text, s
print("PDF sections present: OK")

p = Presentation("docs/tracepulse_ard.pptx")
print("PPTX slides:", len(p.slides))
assert len(p.slides) == 8
titles = [sl.shapes[0].text_frame.text for sl in p.slides]
for t in titles:
    print(" -", t[:60])
ZipFile("docs/tracepulse_ard.pptx").testzip()
print("PPTX zip integrity: OK")
